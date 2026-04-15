#!/usr/bin/env python3
"""
Extract features and images from a Digital Book PowerPoint file and generate
JSON feature files for the digitalbook Next.js app.

Usage:
    python3 scripts/pptx-to-features.py <pptx-path> <edition-slug>

Example:
    python3 scripts/pptx-to-features.py \
        content/editions/april-2026/powerpoint/Digital\ Book\ -\ Apriltemplate.pptx \
        mar-2026

What it does:
    1. Parses all slides, extracting text and embedded images
    2. Detects section boundaries ("In a nutshell" slides)
    3. Generates one JSON feature file per feature slide
    4. Extracts images to public/images/<edition>/<section>/
    5. Outputs a summary of what was created

Prerequisites:
    pip install python-pptx
"""

from __future__ import annotations

import json
import os
import re
import sys
import hashlib
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Section slug mapping — matches the digitalbook content structure
SECTION_KEYWORDS = {
    "web": "website-rc",
    "website": "website-rc",
    "app": "club-avolta-app",
    "oms": "oms",
    "sso": "sso",
    "my autogrill": "my-autogrill",
    "autogrill": "my-autogrill",
}

# Image directory names per section (matches existing conventions)
SECTION_IMAGE_DIRS = {
    "website-rc": "web",
    "club-avolta-app": "App",
    "oms": "OMS",
    "sso": "SSO",
    "my-autogrill": "myautogrill",
}

# Platform tags extracted from shape text
KNOWN_PLATFORMS = [
    "Website", "Reserve & Collect", "Emporium", "Staff Shopping",
    "App", "Club Avolta", "OMS", "SSO", "My Autogrill",
]


def slugify(text: str) -> str:
    """Convert text to a URL-friendly slug."""
    text = text.lower().strip()
    text = re.sub(r'[^\w\s-]', '', text)
    text = re.sub(r'[\s_]+', '-', text)
    text = re.sub(r'-+', '-', text)
    return text.strip('-')[:60]


def extract_slide_text(slide) -> list[str]:
    """Get all non-empty text lines from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def extract_slide_images(slide, output_dir: str, slide_num: int, seen_hashes: set) -> list[str]:
    """Extract images from a slide, returning list of saved file paths."""
    saved = []
    img_idx = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (hasattr(shape, 'image') and shape.image):
            try:
                img = shape.image
                blob = img.blob
                h = hashlib.md5(blob).hexdigest()
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)

                ext = img.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'

                img_idx += 1
                filename = f"slide{slide_num}-{img_idx}.{ext}"
                filepath = os.path.join(output_dir, filename)
                os.makedirs(output_dir, exist_ok=True)
                with open(filepath, 'wb') as f:
                    f.write(blob)
                saved.append(filepath)
            except Exception:
                pass
    return saved


def detect_platforms(texts: list[str]) -> list[str]:
    """Detect platform tags from slide text."""
    all_text = " ".join(texts)
    return [p for p in KNOWN_PLATFORMS if p in all_text]


def detect_section_from_nutshell(texts: list[str]) -> str | None:
    """Detect which section an 'In a nutshell' slide belongs to."""
    all_text = " ".join(texts).lower()
    if "in a nutshell" not in all_text:
        return None

    for keyword, section_slug in SECTION_KEYWORDS.items():
        if keyword in all_text:
            return section_slug
    return None


def is_feature_slide(texts: list[str]) -> bool:
    """Determine if a slide contains a feature (has Goal section or meaningful content)."""
    all_text = " ".join(texts).lower()
    if "in a nutshell" in all_text:
        return False
    if "purpose of" in all_text and "document" in all_text:
        return False
    if "contents" in all_text and len(texts) < 3:
        return False
    if "thank" in all_text and "you" in all_text and len(texts) <= 2:
        return False
    if "digital book" in all_text and len(texts) <= 5 and any("release" in t.lower() for t in texts):
        return False
    # Must have some substantial text
    return len(texts) >= 2


def extract_goal(texts: list[str]) -> str:
    """Extract the goal text from slide content."""
    goal_lines = []
    capture = False
    for t in texts:
        if t.lower().strip() == "goal":
            capture = True
            continue
        if capture:
            # Stop at known section markers
            if t.strip() == "\u200b" or t.strip() == "":
                continue
            goal_lines.append(t)
            if len(goal_lines) >= 2:
                break
    return goal_lines[-1] if goal_lines else ""


def extract_description(texts: list[str]) -> str:
    """Extract the description (non-title, non-goal, non-platform text)."""
    skip = {"goal", "\u200b", ""}
    desc_lines = []
    past_goal = False

    for t in texts[1:]:  # Skip title
        lower = t.lower().strip()
        if lower in skip:
            continue
        if lower == "goal":
            past_goal = True
            continue
        if any(t.strip() == p for p in KNOWN_PLATFORMS):
            continue
        if past_goal:
            # First line after goal section is the goal itself, skip it
            past_goal = False
            continue
        desc_lines.append(t)

    return " ".join(desc_lines) if desc_lines else ""


def format_release_date(edition_slug: str) -> str:
    """Convert 'mar-2026' to 'Mar. 2026'."""
    month_abbr, year = edition_slug.split("-")
    return f"{month_abbr.capitalize()}. {year}"


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 scripts/pptx-to-features.py <pptx-path> <edition-slug>")
        print("Example: python3 scripts/pptx-to-features.py path/to/file.pptx mar-2026")
        sys.exit(1)

    pptx_path = sys.argv[1]
    edition_slug = sys.argv[2]

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    features_base = os.path.join(project_root, "content", "editions", edition_slug)
    images_base = os.path.join(project_root, "public", "images", edition_slug)
    release_date = format_release_date(edition_slug)

    if not os.path.exists(features_base):
        print(f"Error: Edition directory not found: {features_base}")
        print(f"Run 'npm run new-edition {edition_slug}' first.")
        sys.exit(1)

    prs = Presentation(pptx_path)
    print(f"\n📖 Parsing: {pptx_path}")
    print(f"   Edition: {edition_slug}")
    print(f"   Slides: {len(prs.slides)}\n")

    current_section = None
    feature_counts = {}
    seen_hashes = set()
    total_features = 0
    total_images = 0

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        texts = extract_slide_text(slide)

        if not texts:
            continue

        # Check if this is a section boundary slide
        detected = detect_section_from_nutshell(texts)
        if detected:
            current_section = detected
            feature_counts.setdefault(current_section, 0)
            print(f"  📂 Section: {current_section}")
            continue

        # Skip slides before first section, or non-feature slides
        if not current_section:
            continue

        if not is_feature_slide(texts):
            continue

        # This is a feature slide
        feature_counts[current_section] = feature_counts.get(current_section, 0) + 1
        num = feature_counts[current_section]
        title = texts[0]

        # Clean up numbering from title (e.g., "2. UK - Price per unit" -> "UK - Price per unit")
        clean_title = re.sub(r'^\d+\.\s*', '', title)

        slug = slugify(clean_title)
        goal = extract_goal(texts)
        description = extract_description(texts)
        platforms = detect_platforms(texts)

        # Extract images
        img_dir = SECTION_IMAGE_DIRS.get(current_section, current_section)
        img_output_dir = os.path.join(images_base, img_dir)
        saved_images = extract_slide_images(slide, img_output_dir, slide_num, seen_hashes)
        total_images += len(saved_images)

        # Convert to public URL paths
        image_urls = [
            f"/images/{edition_slug}/{img_dir}/{os.path.basename(p)}"
            for p in saved_images
        ]

        # Build feature JSON
        feature = {
            "slug": slug,
            "title": clean_title,
            "status": "go-live",
            "releaseDate": release_date,
            "goal": goal,
            "description": description,
            "tags": [],
            "regions": ["Global"],
            "platforms": platforms if platforms else [],
            "images": image_urls,
        }

        # Write feature file
        features_dir = os.path.join(features_base, current_section, "features")
        os.makedirs(features_dir, exist_ok=True)
        filename = f"{num:02d}-{slug}.json"
        filepath = os.path.join(features_dir, filename)

        with open(filepath, 'w') as f:
            json.dump(feature, f, indent=2, ensure_ascii=False)
            f.write('\n')

        total_features += 1
        img_info = f" + {len(saved_images)} images" if saved_images else ""
        print(f"    ✅ {filename}{img_info}")

    print(f"\n🎉 Done! Created {total_features} features, extracted {total_images} images.")
    print(f"\nNext steps:")
    print(f"  1. Review generated JSON files in content/editions/{edition_slug}/")
    print(f"  2. Fill in missing tags, regions, and descriptions")
    print(f"  3. Run: npm run validate")
    print(f"  4. Run: npm run build\n")


if __name__ == "__main__":
    main()
